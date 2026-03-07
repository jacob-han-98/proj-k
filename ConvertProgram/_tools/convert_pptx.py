"""
PPTX → AI-Readable Markdown 변환 파이프라인
============================================
파워포인트의 슬라이드 텍스트, 노트, 도형을 추출하여 Markdown으로 변환

사용법:
  python convert_pptx.py "경로/파일.pptx"
  python convert_pptx.py --batch "폴더경로"

의존성: python-pptx (pip install python-pptx)
"""

import os
import sys
import json
import re
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("python-pptx가 필요합니다: pip install python-pptx")
    sys.exit(1)


def extract_shape_text(shape):
    """도형에서 텍스트 추출"""
    texts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                # 들여쓰기 레벨 반영
                level = para.level or 0
                prefix = '  ' * level + '- ' if level > 0 else ''
                texts.append(f"{prefix}{text}")
    return texts


def extract_table(shape):
    """테이블 도형에서 Markdown 테이블 추출"""
    if not shape.has_table:
        return ""

    table = shape.table
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip().replace('\n', ' / ').replace('|', '∣'))
        rows.append(cells)

    if not rows:
        return ""

    md = f"| {' | '.join(rows[0])} |\n"
    md += f"| {' | '.join(['---'] * len(rows[0]))} |\n"
    for row in rows[1:]:
        # 열 수 맞추기
        while len(row) < len(rows[0]):
            row.append('')
        md += f"| {' | '.join(row[:len(rows[0])])} |\n"

    return md


def convert_pptx(pptx_path, out_dir=None):
    """단일 PPTX 파일 변환"""
    pptx_path = os.path.abspath(pptx_path)
    filename = os.path.basename(pptx_path)
    stem = Path(pptx_path).stem

    if out_dir is None:
        base = os.path.dirname(pptx_path)
        kb_dir = os.path.join(base, '..', '_knowledge_base', 'pptx')
        out_dir = os.path.abspath(kb_dir)

    os.makedirs(out_dir, exist_ok=True)

    print(f"  변환: {filename}")

    prs = Presentation(pptx_path)

    md = f"# {stem}\n\n"
    md += f"> 원본: {filename}\n"
    md += f"> 슬라이드: {len(prs.slides)}장\n\n"

    for i, slide in enumerate(prs.slides, 1):
        # 슬라이드 제목
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        md += f"---\n\n## 슬라이드 {i}"
        if title:
            md += f": {title}"
        md += "\n\n"

        # 모든 도형 처리
        for shape in slide.shapes:
            # 테이블
            if shape.has_table:
                md += extract_table(shape) + "\n"
            # 텍스트 도형
            elif shape.has_text_frame:
                # 제목은 이미 처리됨
                if shape == slide.shapes.title:
                    continue
                texts = extract_shape_text(shape)
                if texts:
                    md += '\n'.join(texts) + '\n\n'
            # 그룹 도형
            elif shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                try:
                    for child in shape.shapes:
                        if hasattr(child, 'has_text_frame') and child.has_text_frame:
                            texts = extract_shape_text(child)
                            if texts:
                                md += '\n'.join(texts) + '\n'
                except:
                    pass

        # 슬라이드 노트
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                md += f"\n> **노트:** {notes}\n\n"

    safe_name = stem.replace(' ', '_')
    out_path = os.path.join(out_dir, f"{safe_name}.md")
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(md)

    print(f"    → {len(md):,} chars, {len(prs.slides)} slides")

    return {
        'status': 'success',
        'filename': filename,
        'output': out_path,
        'slides': len(prs.slides),
        'chars': len(md)
    }


def main():
    parser = argparse.ArgumentParser(description='PPTX → AI-Readable Markdown 변환')
    parser.add_argument('input', nargs='?', help='PPTX 파일 경로')
    parser.add_argument('--out', '-o', help='출력 폴더')
    parser.add_argument('--batch', '-b', help='폴더 내 모든 PPTX 일괄 변환')

    args = parser.parse_args()

    if args.batch:
        batch_dir = os.path.abspath(args.batch)
        pptx_files = []
        for root, dirs, files in os.walk(batch_dir):
            dirs[:] = [d for d in dirs if d != '_knowledge_base' and d != '_tools']
            for f in files:
                if f.endswith('.pptx') and not f.startswith('~$'):
                    pptx_files.append(os.path.join(root, f))

        print(f"발견된 PPTX 파일: {len(pptx_files)}개\n")

        out_dir = args.out or os.path.join(batch_dir, '_knowledge_base', 'pptx')
        os.makedirs(out_dir, exist_ok=True)

        results = {}
        success = 0
        for i, pptx in enumerate(sorted(pptx_files), 1):
            print(f"[{i}/{len(pptx_files)}]", end='')
            try:
                r = convert_pptx(pptx, out_dir)
                results[pptx] = r
                if r['status'] == 'success':
                    success += 1
            except Exception as e:
                print(f"  오류: {e}")
                results[pptx] = {'status': 'error', 'error': str(e)}

        index_path = os.path.join(out_dir, "_INDEX.json")
        with open(index_path, 'w', encoding='utf-8') as f:
            json.dump({
                'total': len(results),
                'success': success,
                'total_chars': sum(r.get('chars', 0) for r in results.values()),
                'files': {os.path.basename(k): v for k, v in results.items()}
            }, f, ensure_ascii=False, indent=2)

        print(f"\n{'='*60}")
        print(f"배치 완료: {success}/{len(results)} 성공")
        return results

    elif args.input:
        return convert_pptx(args.input, args.out)
    else:
        parser.print_help()


if __name__ == '__main__':
    main()
